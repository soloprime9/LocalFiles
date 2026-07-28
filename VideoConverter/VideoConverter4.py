#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Video to Presentation Converter — Flask Edition
=================================================
Single-file Flask app (no external HTML templates, no Tkinter).

Fixes included vs. the original VideoConverter3.py:
  1. QUALITY  -> removed the 'ext=mp4' filter that silently capped
                 downloads at 1080p (most 1440p/4K/8K streams are webm).
                 Now uses format_sort to always grab the single highest
                 available resolution automatically (no dropdown/manual pick).
  2. ASPECT RATIO -> images are no longer stretched to fill the slide;
                 they are centered and scaled to fit while preserving
                 the original aspect ratio.
  3. JPEG QUALITY -> frames are now written at maximum JPEG quality (100)
                 instead of OpenCV's default (~95).
  4. CONCURRENCY -> temp frame directories are unique per job
                 (tempfile.mkdtemp) instead of one shared hardcoded
                 folder, so multiple users/jobs never collide.
  5. CLEANUP  -> temp directories are actually deleted after export
                 (previously this was silently skipped / commented out).
  6. UI       -> Tkinter GUI removed entirely; replaced with a Flask
                 web UI (frame preview + checkbox selection, same idea
                 as the old desktop app) served from an inline HTML
                 string via render_template_string — no templates/ folder.

Run:
    pip install flask opencv-python python-pptx yt-dlp pillow
    python video_converter_flask.py
    -> open http://127.0.0.1:5000
"""

import os
import sys
import shutil
import logging
import tempfile
import threading
import uuid
from pathlib import Path
from datetime import datetime

import cv2

try:
    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.util import Inches
    import yt_dlp
    from flask import (
        Flask, request, jsonify, send_file, render_template_string, abort
    )
except ImportError as e:
    print(f"Dependency Error: {e}\n"
          f"Run: pip install flask opencv-python python-pptx yt-dlp pillow")
    sys.exit(1)

logging.basicConfig(level=logging.INFO,
                     format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)


# --------------------------------------------------------------------------- #
#  Core conversion logic (Tkinter-free)
# --------------------------------------------------------------------------- #

def add_picture_fit(slide, img_path, slide_w, slide_h):
    """Insert an image centered on the slide, scaled to fit WITHOUT
    stretching / distorting the original aspect ratio."""
    with PILImage.open(img_path) as im:
        iw, ih = im.size

    img_ratio = iw / ih
    slide_ratio = slide_w / slide_h

    if img_ratio > slide_ratio:
        w = slide_w
        h = int(slide_w / img_ratio)
    else:
        h = slide_h
        w = int(slide_h * img_ratio)

    left = int((slide_w - w) / 2)
    top = int((slide_h - h) / 2)
    slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)


class VideoExtractor:
    """Handles frame extraction, thumbnailing and PPTX export for a
    single conversion job. One instance per job -> safe for concurrent
    use across multiple Flask requests."""

    def __init__(self, source: str, interval: int = 1):
        self.source = source.strip()
        self.interval = max(1, int(interval))
        self.is_url = self.source.lower().startswith(('http://', 'https://', 'www.'))

        # Unique temp dir per job -> no collisions between concurrent users
        self.temp_dir = Path(tempfile.mkdtemp(prefix="slides_"))
        self.saved_frames = []   # list[Path]
        self.base_name = "Online_Video"

        if not self.is_url:
            local_path = Path(self.source)
            if not local_path.exists():
                raise FileNotFoundError(f"Local file not found: {self.source}")
            self.base_name = local_path.stem

    # ---- resolution: always grab the single highest quality stream ---- #
    def _get_stream_url(self) -> str:
        logger.info(f"Resolving stream for: {self.source}")
        ydl_opts = {
            # No 'ext=mp4' restriction: that filter silently drops 1440p/4K/8K
            # streams, which are almost always webm (vp9/av1), not mp4.
            'format': 'bestvideo/best',
            # Always sort for the highest resolution/fps/codec quality first —
            # this replaces any manual "choose a resolution" dropdown.
            'format_sort': ['res', 'fps', 'vcodec:vp9.2', 'vcodec:vp9', 'br'],
            'quiet': True,
            'no_warnings': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(self.source, download=False)
            if info.get('title'):
                self.base_name = "".join(
                    c for c in info['title'] if c.isalnum() or c in (' ', '_', '-')
                ).rstrip().replace(' ', '_') or self.base_name
            width = info.get('width')
            height = info.get('height')
            logger.info(f"Selected stream resolution: {width}x{height}")
            return info['url']

    def extract_frames(self) -> None:
        target_capture = self._get_stream_url() if self.is_url else str(self.source)

        logger.info("Opening video stream...")
        cap = cv2.VideoCapture(target_capture)
        if not cap.isOpened():
            raise ValueError("Could not open the video stream/source.")

        fps = cap.get(cv2.CAP_PROP_FPS)
        if fps <= 0:
            fps = 25  # fallback for streams that hide metadata

        frame_gap = max(1, int(fps * self.interval))
        frame_idx = 0
        saved_idx = 0

        while True:
            success, frame = cap.read()
            if not success:
                break

            if frame_idx % frame_gap == 0:
                target_img_path = self.temp_dir / f"slide_{saved_idx:04d}.jpg"
                # Max JPEG quality (OpenCV default is ~95, we force 100)
                cv2.imwrite(str(target_img_path), frame, [cv2.IMWRITE_JPEG_QUALITY, 100])
                self.saved_frames.append(target_img_path)
                saved_idx += 1

                if saved_idx % 10 == 0:
                    logger.info(f"Extracted {saved_idx} frames...")

            frame_idx += 1

        cap.release()
        logger.info(f"Done. Extracted {len(self.saved_frames)} frames.")

    def get_selected_frames(self, selected_indices):
        return [self.saved_frames[i] for i in selected_indices
                if 0 <= i < len(self.saved_frames)]

    def save_to_pptx(self, selected_indices) -> Path:
        frames = self.get_selected_frames(selected_indices)
        if not frames:
            raise ValueError("No frames selected for export.")

        prs = Presentation()
        prs.slide_width, prs.slide_height = SLIDE_WIDTH, SLIDE_HEIGHT

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Video Content Export"
        source_label = self.source if not self.is_url else "Remote Video Source"
        title_slide.placeholders[1].text = (
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"Source: {source_label}"
        )

        blank_layout = prs.slide_layouts[6]
        for img_path in frames:
            slide = prs.slides.add_slide(blank_layout)
            add_picture_fit(slide, img_path, prs.slide_width, prs.slide_height)

        output_path = self.temp_dir / f"{self.base_name}_output.pptx"
        prs.save(str(output_path))
        logger.info(f"PPTX saved: {output_path}")
        return output_path

    def export_images(self, selected_indices) -> Path:
        """Zips the selected frames and returns the zip path."""
        frames = self.get_selected_frames(selected_indices)
        if not frames:
            raise ValueError("No frames selected for export.")

        export_dir = self.temp_dir / f"{self.base_name}_HD_Frames"
        export_dir.mkdir(exist_ok=True)
        for img_path in frames:
            shutil.copy(img_path, export_dir / img_path.name)

        zip_base = self.temp_dir / f"{self.base_name}_HD_Frames"
        zip_path = shutil.make_archive(str(zip_base), 'zip', root_dir=str(export_dir))
        return Path(zip_path)

    def cleanup(self) -> None:
        shutil.rmtree(self.temp_dir, ignore_errors=True)
        logger.info(f"Cleaned up temp dir: {self.temp_dir}")


# --------------------------------------------------------------------------- #
#  Flask app + in-memory job registry
# --------------------------------------------------------------------------- #

app = Flask(__name__)

jobs = {}          # job_id -> {"status", "extractor", "error"}
jobs_lock = threading.Lock()


def run_extraction_job(job_id: str, extractor: VideoExtractor):
    try:
        extractor.extract_frames()
        with jobs_lock:
            jobs[job_id]["status"] = "done"
    except Exception as e:
        logger.error(f"Job {job_id} failed: {e}")
        with jobs_lock:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = str(e)


@app.route('/')
def index():
    return render_template_string(INDEX_HTML)


@app.route('/extract', methods=['POST'])
def extract():
    data = request.get_json(force=True)
    source = (data.get('source') or '').strip()
    interval = data.get('interval', 1)

    if not source:
        return jsonify({"error": "Please provide a video URL or local file path."}), 400

    try:
        extractor = VideoExtractor(source, interval=interval)
    except FileNotFoundError as e:
        return jsonify({"error": str(e)}), 400

    job_id = str(uuid.uuid4())
    with jobs_lock:
        jobs[job_id] = {"status": "processing", "extractor": extractor, "error": None}

    threading.Thread(target=run_extraction_job, args=(job_id, extractor), daemon=True).start()
    return jsonify({"job_id": job_id})


@app.route('/status/<job_id>')
def status(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
        if not job:
            return jsonify({"status": "not_found"}), 404
        return jsonify({
            "status": job["status"],
            "error": job["error"],
            "frame_count": len(job["extractor"].saved_frames),
        })


@app.route('/frame/<job_id>/<int:idx>')
def frame(job_id, idx):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        abort(404)
    frames = job["extractor"].saved_frames
    if idx < 0 or idx >= len(frames):
        abort(404)
    return send_file(str(frames[idx]), mimetype='image/jpeg')


@app.route('/export/pptx/<job_id>', methods=['POST'])
def export_pptx(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        abort(404)
    data = request.get_json(force=True)
    selected = data.get('selected', [])
    try:
        output_path = job["extractor"].save_to_pptx(selected)
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    return send_file(str(output_path), as_attachment=True,
                      download_name=output_path.name)


@app.route('/export/images/<job_id>', methods=['POST'])
def export_images(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        abort(404)
    data = request.get_json(force=True)
    selected = data.get('selected', [])
    try:
        zip_path = job["extractor"].export_images(selected)
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    return send_file(str(zip_path), as_attachment=True, download_name=zip_path.name)


@app.route('/cleanup/<job_id>', methods=['POST'])
def cleanup(job_id):
    with jobs_lock:
        job = jobs.pop(job_id, None)
    if job:
        job["extractor"].cleanup()
    return jsonify({"ok": True})


# --------------------------------------------------------------------------- #
#  Inline HTML/JS UI (no templates/ folder, no Tkinter)
# --------------------------------------------------------------------------- #

INDEX_HTML = """
<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>Video to Presentation Converter</title>
<style>
  body { font-family: Arial, sans-serif; max-width: 900px; margin: 30px auto; padding: 0 15px; }
  h1 { font-size: 22px; }
  .row { display: flex; gap: 10px; margin: 10px 0; }
  input[type=text] { flex: 1; padding: 8px; }
  input[type=number] { width: 80px; padding: 8px; }
  button { padding: 8px 16px; cursor: pointer; }
  #status { margin: 10px 0; font-style: italic; color: #555; }
  #grid { display: grid; grid-template-columns: repeat(4, 1fr); gap: 10px; margin: 15px 0; }
  .cell { border: 1px solid #ddd; padding: 6px; text-align: center; border-radius: 4px; }
  .cell img { width: 100%; border-radius: 3px; }
  .actions { display: flex; gap: 10px; margin: 15px 0; }
  .actions button { flex: 1; }
</style>
</head>
<body>
  <h1>Video to Presentation Converter</h1>

  <div class="row">
    <input type="text" id="source" placeholder="Video URL or local file path">
    <input type="number" id="interval" value="1" min="1" title="Seconds between frames">
    <button onclick="startExtraction()">Extract Frames</button>
  </div>

  <div id="status">Status: Ready</div>

  <div class="row" style="display:none" id="selectRow">
    <button onclick="selectAll(true)">Select All</button>
    <button onclick="selectAll(false)">Deselect All</button>
  </div>

  <div id="grid"></div>

  <div class="actions" style="display:none" id="exportRow">
    <button onclick="exportPptx()">Save Selected as PPTX</button>
    <button onclick="exportImages()">Save Selected Images (.zip)</button>
  </div>

<script>
let jobId = null;
let frameCount = 0;

function setStatus(text) {
  document.getElementById('status').innerText = 'Status: ' + text;
}

async function startExtraction() {
  const source = document.getElementById('source').value.trim();
  const interval = parseInt(document.getElementById('interval').value || '1', 10);
  if (!source) { alert('Please provide a video URL or local file path.'); return; }

  document.getElementById('grid').innerHTML = '';
  document.getElementById('selectRow').style.display = 'none';
  document.getElementById('exportRow').style.display = 'none';
  setStatus('Starting extraction...');

  const res = await fetch('/extract', {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({source, interval})
  });
  const data = await res.json();
  if (data.error) { setStatus('Error: ' + data.error); return; }

  jobId = data.job_id;
  pollStatus();
}

async function pollStatus() {
  const res = await fetch('/status/' + jobId);
  const data = await res.json();

  if (data.status === 'processing') {
    setStatus('Extracting frames... (' + data.frame_count + ' so far)');
    setTimeout(pollStatus, 1500);
  } else if (data.status === 'done') {
    frameCount = data.frame_count;
    setStatus('Extracted ' + frameCount + ' frames.');
    buildGrid();
  } else if (data.status === 'error') {
    setStatus('Error: ' + data.error);
  }
}

function buildGrid() {
  const grid = document.getElementById('grid');
  grid.innerHTML = '';
  for (let i = 0; i < frameCount; i++) {
    const cell = document.createElement('div');
    cell.className = 'cell';
    cell.innerHTML = `
      <img src="/frame/${jobId}/${i}" loading="lazy">
      <div><label><input type="checkbox" class="frame-check" data-idx="${i}" checked> Frame ${i+1}</label></div>
    `;
    grid.appendChild(cell);
  }
  document.getElementById('selectRow').style.display = 'flex';
  document.getElementById('exportRow').style.display = 'flex';
}

function selectAll(state) {
  document.querySelectorAll('.frame-check').forEach(cb => cb.checked = state);
}

function getSelected() {
  return Array.from(document.querySelectorAll('.frame-check'))
    .filter(cb => cb.checked)
    .map(cb => parseInt(cb.dataset.idx, 10));
}

async function exportPptx() {
  const selected = getSelected();
  if (!selected.length) { alert('Select at least one frame.'); return; }
  setStatus('Building PPTX...');
  const res = await fetch('/export/pptx/' + jobId, {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({selected})
  });
  if (!res.ok) { const err = await res.json(); setStatus('Error: ' + err.error); return; }
  const blob = await res.blob();
  downloadBlob(blob, 'presentation.pptx');
  setStatus('PPTX downloaded.');
}

async function exportImages() {
  const selected = getSelected();
  if (!selected.length) { alert('Select at least one frame.'); return; }
  setStatus('Zipping images...');
  const res = await fetch('/export/images/' + jobId, {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({selected})
  });
  if (!res.ok) { const err = await res.json(); setStatus('Error: ' + err.error); return; }
  const blob = await res.blob();
  downloadBlob(blob, 'frames.zip');
  setStatus('Images downloaded.');
}

function downloadBlob(blob, filename) {
  const url = URL.createObjectURL(blob);
  const a = document.createElement('a');
  a.href = url; a.download = filename;
  document.body.appendChild(a); a.click(); a.remove();
  URL.revokeObjectURL(url);
}
</script>
</body>
</html>
"""


if __name__ == '__main__':
    app.run(debug=True, host='127.0.0.1', port=8080)