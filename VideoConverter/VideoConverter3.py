#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import cv2
import shutil
import logging
import threading
from pathlib import Path
from datetime import datetime

import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from PIL import Image, ImageTk  # Added for displaying thumbnails safely

try:
    from pptx import Presentation
    from pptx.util import Inches
    import yt_dlp
except ImportError as e:
    print(f"Dependency Error: {e}\nRun: pip install opencv-python python-pptx yt-dlp pillow")
    sys.exit(1)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class ModernVideoConverterGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("HD Video Frame & Presentation Exporter")
        self.root.geometry("750x650")  # Expanded to fit the preview panel comfortably
        self.root.minimum_size = (700, 550)
        
        self.temp_dir = Path("extracted_slides_cache")
        self.saved_frames = []      # List of Path objects
        self.frame_vars = []        # Holds BooleanVars linked to checkboxes
        self.thumbnail_images = []  # Keeps Python references so GC doesn't wipe them
        self.base_name = "Online_Video"
        self.is_url = False
        
        # Style
        self.style = ttk.Style()
        self.style.theme_use('clam')
        
        self.create_widgets()

    def create_widgets(self):
        # Main Container
        main_frame = ttk.Frame(self.root, padding="15")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # Header text
        header = ttk.Label(main_frame, text="Video to HD Image & Presentation Converter", font=("Arial", 14, "bold"))
        header.pack(pady=(0, 10))
        
        # 1. Input Source Selection Row
        source_frame = ttk.LabelFrame(main_frame, text=" 1. Video Source (URL or Local File Path) ", padding="10")
        source_frame.pack(fill=tk.X, pady=5)
        
        self.source_entry = ttk.Entry(source_frame, width=50)
        self.source_entry.pack(side=tk.LEFT, padx=(0, 10), expand=True, fill=tk.X)
        
        browse_btn = ttk.Button(source_frame, text="Browse File", command=self.browse_local_file)
        browse_btn.pack(side=tk.RIGHT)
        
        # Configuration & Extraction Control Row
        config_frame = ttk.Frame(main_frame, padding="5")
        config_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(config_frame, text="Extract Frame Every:").pack(side=tk.LEFT, padx=(0, 5))
        self.interval_spin = ttk.Spinbox(config_frame, from_=1, to=60, width=5)
        self.interval_spin.set(1)
        self.interval_spin.pack(side=tk.LEFT, padx=(0, 5))
        ttk.Label(config_frame, text="seconds").pack(side=tk.LEFT)
        
        self.process_btn = ttk.Button(config_frame, text="⚡ Extract HD Frames", command=self.start_processing_thread)
        self.process_btn.pack(side=tk.RIGHT, padx=5)

        # Status box
        self.status_label = ttk.Label(main_frame, text="Status: Ready", font=("Arial", 10, "italic"), foreground="blue")
        self.status_label.pack(pady=5)
        
        # 2. Frame Preview & Selection UI Section
        preview_label_frame = ttk.LabelFrame(main_frame, text=" 2. Preview & Select Frames to Keep ", padding="10")
        preview_label_frame.pack(fill=tk.BOTH, expand=True, pady=5)
        
        # Quick Selection Utility Buttons
        select_action_frame = ttk.Frame(preview_label_frame)
        select_action_frame.pack(fill=tk.X, pady=(0, 5))
        
        self.select_all_btn = ttk.Button(select_action_frame, text="✓ Select All", command=self.select_all_frames, state=tk.DISABLED)
        self.select_all_btn.pack(side=tk.LEFT, padx=2)
        self.deselect_all_btn = ttk.Button(select_action_frame, text="✗ Deselect All", command=self.deselect_all_frames, state=tk.DISABLED)
        self.deselect_all_btn.pack(side=tk.LEFT, padx=2)
        
        # Scrollable Canvas Framework for handling thumbnails dynamic rows
        self.canvas = tk.Canvas(preview_label_frame, borderwidth=0, background="#ffffff")
        self.scrollbar = ttk.Scrollbar(preview_label_frame, orient="vertical", command=self.canvas.yview)
        self.scrollable_frame = tk.Frame(self.canvas, background="#ffffff")
        
        self.scrollable_frame.bind(
            "<Configure>",
            lambda e: self.canvas.configure(scrollregion=self.canvas.bbox("all"))
        )
        self.canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        self.canvas.configure(yscrollcommand=self.scrollbar.set)
        
        self.canvas.pack(side="left", fill="both", expand=True)
        self.scrollbar.pack(side="right", fill="y")
        
        # 3. Save / Export Actions Panel Row
        self.btn_frame = ttk.LabelFrame(main_frame, text=" 3. Save Selected Content ", padding="15")
        self.btn_frame.pack(fill=tk.X, pady=(10, 0))
        
        self.save_folder_btn = ttk.Button(self.btn_frame, text="📁 Save Selected Images to Folder", state=tk.DISABLED, command=self.save_images_to_custom_folder)
        self.save_folder_btn.pack(side=tk.LEFT, padx=5, expand=True, fill=tk.X)
        
        self.save_pptx_btn = ttk.Button(self.btn_frame, text="📊 Save Selected as Widescreen PPTX", state=tk.DISABLED, command=self.save_to_pptx)
        self.save_pptx_btn.pack(side=tk.LEFT, padx=5, expand=True, fill=tk.X)

    def browse_local_file(self):
        file_path = filedialog.askopenfilename(filetypes=[("Video Files", "*.mp4 *.mkv *.avi *.mov")])
        if file_path:
            self.source_entry.delete(0, tk.END)
            self.source_entry.insert(0, file_path)

    def update_status(self, text, color="black"):
        self.status_label.config(text=f"Status: {text}", foreground=color)
        self.root.update_idletasks()

    def start_processing_thread(self):
        source = self.source_entry.get().strip()
        if not source:
            messagebox.showerror("Error", "Please provide a YouTube URL or browse for a local video file.")
            return
            
        self.process_btn.config(state=tk.DISABLED)
        self.save_folder_btn.config(state=tk.DISABLED)
        self.save_pptx_btn.config(state=tk.DISABLED)
        self.select_all_btn.config(state=tk.DISABLED)
        self.deselect_all_btn.config(state=tk.DISABLED)
        
        # Completely clear old thumbnails display UI structure
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
            
        threading.Thread(target=self.extract_frames_core, args=(source,), daemon=True).start()

    def _get_stream_url(self, source) -> str:
        ydl_opts = {
            'format': 'bestvideo[ext=mp4]/best[ext=mp4]/best',
            'quiet': True,
            'no_warnings': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(source, download=False)
            if 'title' in info and info['title']:
                self.base_name = "".join([c for c in info['title'] if c.isalnum() or c in (' ', '_', '-')]).rstrip().replace(' ', '_')
            return info['url']

    def extract_frames_core(self, source):
        try:
            self.saved_frames = []
            self.frame_vars = []
            self.thumbnail_images = []
            self.is_url = source.lower().startswith(('http://', 'https://', 'www.'))
            interval = int(self.interval_spin.get())
            
            if self.is_url:
                self.update_status("Contacting online video host stream...", "orange")
                target_capture = self._get_stream_url(source)
            else:
                local_path = Path(source)
                if not local_path.exists():
                    raise FileNotFoundError("Local file does not exist.")
                self.base_name = local_path.stem
                target_capture = str(source)
                
            self.update_status("Opening high-definition media stream...", "orange")
            
            if self.temp_dir.exists():
                shutil.rmtree(self.temp_dir)
            self.temp_dir.mkdir(exist_ok=True)
            
            cap = cv2.VideoCapture(target_capture)
            if not cap.isOpened():
                raise ValueError("Could not extract frames from the target video stream.")
                
            fps = cap.get(cv2.CAP_PROP_FPS)
            if fps <= 0:
                fps = 25
                
            frame_gap = max(1, int(fps * interval))
            frame_idx = 0
            saved_idx = 0
            
            self.update_status("Extracting pristine HD images. Please wait...", "blue")
            
            while True:
                success, frame = cap.read()
                if not success:
                    break
                    
                if frame_idx % frame_gap == 0:
                    target_img_path = self.temp_dir / f"slide_{saved_idx:04d}.jpg"
                    cv2.imwrite(str(target_img_path), frame)
                    self.saved_frames.append(target_img_path)
                    saved_idx += 1
                    
                frame_idx += 1
                
            cap.release()
            
            if self.saved_frames:
                self.update_status(f"Success! Extracted {len(self.saved_frames)} images. Generating preview...", "orange")
                # Trigger Tkinter UI building for image grid inside safe callback scope
                self.root.after(0, self.populate_image_grid)
            else:
                self.update_status("No images captured. Check your intervals.", "red")
                self.process_btn.config(state=tk.NORMAL)
                
        except Exception as e:
            self.update_status(f"Error encountered: {str(e)}", "red")
            messagebox.showerror("Execution Error", f"Something went wrong:\n{str(e)}")
            self.process_btn.config(state=tk.NORMAL)

    def populate_image_grid(self):
        """Builds a responsive visual grid layout showing thumbnails with selection checkboxes."""
        columns = 4  # Set count of thumbnail elements visible horizontally per row
        
        for index, img_path in enumerate(self.saved_frames):
            # Create a localized cell container inside scroll frame view
            cell = ttk.Frame(self.scrollable_frame, padding=5, style="Card.TFrame")
            row = index // columns
            col = index % columns
            cell.grid(row=row, column=col, padx=8, pady=8, sticky="nsew")
            
            try:
                # Open original disk cache file image, downsize sample frame to dynamic icon spec
                pil_img = Image.open(img_path)
                pil_img.thumbnail((140, 105))  # Maintains a cleaner aspect 4:3 standard display bounding box
                tk_thumb = ImageTk.PhotoImage(pil_img)
                self.thumbnail_images.append(tk_thumb)  # Retain critical instance
                
                # Render Graphic Asset object inside localized grid viewport
                img_lbl = tk.Label(cell, image=tk_thumb, bg="#ffffff")
                img_lbl.pack(side=tk.TOP)
            except Exception as io_err:
                err_lbl = ttk.Label(cell, text="[Image Error]")
                err_lbl.pack(side=tk.TOP)
                logger.error(f"Thumbnail load fault matching target frame path asset: {io_err}")

            # Instantiate Checkbox management state tracker object instance
            var = tk.BooleanVar(value=True)  # Check items automatically by default
            self.frame_vars.append(var)
            
            chk = ttk.Checkbutton(cell, text=f"Frame {index+1}", variable=var)
            chk.pack(side=tk.TOP, pady=2)
            
        self.update_status(f"Loaded {len(self.saved_frames)} assets ready for validation storage setup.", "green")
        self.process_btn.config(state=tk.NORMAL)
        self.save_folder_btn.config(state=tk.NORMAL)
        self.save_pptx_btn.config(state=tk.NORMAL)
        self.select_all_btn.config(state=tk.NORMAL)
        self.deselect_all_btn.config(state=tk.NORMAL)

    def select_all_frames(self):
        for var in self.frame_vars:
            var.set(True)

    def deselect_all_frames(self):
        for var in self.frame_vars:
            var.set(False)

    def get_selected_frames(self):
        """Filters the runtime source list tracking exclusively Checked checkboxes."""
        selected = [self.saved_frames[i] for i, var in enumerate(self.frame_vars) if var.get()]
        return selected

    def save_images_to_custom_folder(self):
        selected_frames = self.get_selected_frames()
        if not selected_frames:
            messagebox.showwarning("No Selection", "Please check/select at least one frame image checkbox above.")
            return

        target_destination = filedialog.askdirectory(title="Select Destination Folder to Export HD Images")
        if target_destination:
            dest_path = Path(target_destination) / f"{self.base_name}_HD_Frames"
            dest_path.mkdir(exist_ok=True)
            
            for img_path in selected_frames:
                shutil.copy(img_path, dest_path / img_path.name)
                
            messagebox.showinfo("Export Perfect", f"{len(selected_frames)} selected HD images saved successfully to:\n{dest_path}")

    def save_to_pptx(self):
        selected_frames = self.get_selected_frames()
        if not selected_frames:
            messagebox.showwarning("No Selection", "Please check/select at least one frame image checkbox above.")
            return

        file_dest = filedialog.asksaveasfilename(
            initialfile=f"{self.base_name}_output.pptx",
            filetypes=[("PowerPoint Document", "*.pptx")],
            defaultextension=".pptx"
        )
        if file_dest:
            self.update_status("Creating PPTX presentation asset...", "orange")
            prs = Presentation()
            prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
            
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "Video Content Export"
            title_slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            
            blank_layout = prs.slide_layouts[6]
            for img_path in selected_frames:
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                
            prs.save(file_dest)
            self.update_status("PowerPoint compiled successfully!", "green")
            messagebox.showinfo("Saved", f"Presentation successfully written to:\n{file_dest}")

if __name__ == '__main__':
    root = tk.Tk()
    app = ModernVideoConverterGUI(root)
    root.mainloop()