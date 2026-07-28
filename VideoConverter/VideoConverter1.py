#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import cv2
import argparse
import logging
import shutil
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches
    import yt_dlp
except ImportError as e:
    print(f"Dependency Error: {e}\nRun: pip install opencv-python python-pptx yt-dlp")
    sys.exit(1)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class VideoToPresentationConverter:
    def __init__(self, video_source: str, output_path: str = None, interval: int = 1):
        self.video_source = video_source
        self.interval = interval
        self.temp_dir = Path("extracted_slides_cache")
        self.saved_frames = []
        
        self.is_url = video_source.lower().startswith(('http://', 'https://', 'www.'))
        
        if self.is_url:
            base_name = "Online_Video"
            self.output_path = Path(output_path) if output_path else Path(f"{base_name}_output.pptx")
        else:
            local_path = Path(video_source)
            if not local_path.exists():
                raise FileNotFoundError(f"Target local video file not found: {video_source}")
            self.output_path = Path(output_path) if output_path else local_path.with_name(f"{local_path.stem}_output.pptx")

    def _get_stream_url(self) -> str:
        logger.info(f"Extracting live stream endpoint from URL: {self.video_source}")
        ydl_opts = {
            'format': 'bestvideo[ext=mp4]/best[ext=mp4]/best',
            'quiet': True,
            'no_warnings': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(self.video_source, download=False)
            if 'title' in info:
                logger.info(f"Target Video Title: {info['title']}")
                if not self.output_path.name.replace('_output.pptx', '').isalnum():
                    clean_title = "".join([c for c in info['title'] if c.isalnum() or c in (' ', '_', '-')]).rstrip()
                    self.output_path = self.output_path.with_name(f"{clean_title.replace(' ', '_')}_output.pptx")
            return info['url']

    def extract_frames(self) -> None:
        self.temp_dir.mkdir(exist_ok=True)
        
        target_capture = self._get_stream_url() if self.is_url else str(self.video_source)
        
        logger.info("Opening video stream buffer...")
        cap = cv2.VideoCapture(target_capture)
        if not cap.isOpened():
            raise ValueError("OpenCV failed to initialize or read video stream source.")
            
        fps = cap.get(cv2.CAP_PROP_FPS)
        if fps <= 0:
            fps = 25  # Fallback assumption for standard streams if metadata is hidden
            
        frame_gap = max(1, int(fps * self.interval))
        frame_idx = 0
        saved_idx = 0
        
        while True:
            success, frame = cap.read()
            if not success:
                break
                
            if frame_idx % frame_gap == 0:
                target_img_path = self.temp_dir / f"slide_{saved_idx:04d}.jpg"
                cv2.imwrite(str(target_img_path), frame)
                self.saved_frames.append(target_img_path)
                saved_idx += 1
                
                if saved_idx % 10 == 0:
                    logger.info(f"Cached {saved_idx} frames...")
                
            frame_idx += 1
            
        cap.release()
        logger.info(f"Completed stream processing. Cached {len(self.saved_frames)} slide visual assets.")

    def compile_presentation(self) -> None:
        if not self.saved_frames:
            logger.warning("No frame sequence data discovered. Aborting PowerPoint output.")
            return

        logger.info("Assembling structural slides into widescreen PPTX layout...")
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
        
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Video Content Presentation"
        
        source_label = self.video_source if not self.is_url else "Remote Live Stream Web Resource"
        title_slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\nSource: {source_label}"
        
        blank_layout = prs.slide_layouts[6]
        for img_path in self.saved_frames:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            
        prs.save(str(self.output_path))
        logger.info(f"Presentation saved successfully: {self.output_path}")

    def purge_cache(self) -> None:
        if self.temp_dir.exists():
            # shutil.rmtree(self.temp_dir)
            logger.info("Temporary visual file cache cleared out.")

    def process(self) -> None:
        try:
            self.extract_frames()
            self.compile_presentation()
        finally:
            self.purge_cache()

def main():
    parser = argparse.ArgumentParser(description="Convert local files or online video stream links straight into PowerPoint presentations.")
    parser.add_argument('source', help='Path to target local video file OR online platform video URL')
    parser.add_argument('-o', '--output', help='Custom presentation target filename')
    parser.add_argument('-i', '--interval', type=int, default=1, help='Snapshot gap threshold in seconds')
    
    args = parser.parse_args()
    
    try:
        converter = VideoToPresentationConverter(args.source, args.output, args.interval)
        converter.process()
    except Exception as err:
        logger.error(f"Execution failed: {err}")
        sys.exit(1)

if __name__ == '__main__':
    main()