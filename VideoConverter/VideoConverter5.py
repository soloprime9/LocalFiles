#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Video to Presentation Converter — Flask Edition (v2, advanced)
================================================================
Single-file Flask app. No templates/ folder, no Tkinter — HTML/CSS/JS
lives inline in the INDEX_HTML string below.

What's new in this version (on top of the previous fixes):
  1. LIVE LOG        -> every step (resolving source, detected quality,
                         opening stream, each saved frame) is pushed to
                         the browser as a line-by-line console feed.
  2. PROGRESS BAR     -> real % progress computed from elapsed video time
                         vs. total duration (falls back to a frame
                         counter for live streams with unknown length).
  3. QUALITY BADGE    -> resolution is detected up front (via yt-dlp for
                         URLs, via OpenCV for local files) and shown as
                         a readable label, e.g. "1920x1080 — Full HD".
  4. LIVE THUMBNAILS  -> the frame grid fills in as frames are produced,
                         not after the whole extraction finishes.
  5. LIGHTBOX/CAROUSEL-> clicking any thumbnail opens a full-size viewer
                         with Prev/Next (and arrow-key) navigation.
  6. ON-DISK PREVIEW  -> frames are written directly into a
                         `preview_frames/` folder next to this script
                         (browsable in your normal file explorer). It is
                         wiped clean at the start of every new
                         extraction — nothing accumulates, and nothing
                         is "kept" unless you explicitly export it.
  7. NO FULL DOWNLOAD -> remote sources are never downloaded to disk in
                         full; yt-dlp resolves the highest-quality
                         *stream URL* and OpenCV reads frames from that
                         stream directly.
  8. UI               -> gradient background, glass cards, animated
                         progress bar, spinner, fade-in thumbnails.

Run:
    pip install flask opencv-python python-pptx yt-dlp pillow
    python video_converter_flask.py
    -> open http://127.0.0.1:5000
"""

import os
import sys
import shutil
import logging
import tempfile
import threading
import uuid
from pathlib import Path
from datetime import datetime

import cv2

try:
    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.util import Inches
    import yt_dlp
    from flask import (
        Flask, request, jsonify, send_file, render_template_string, abort
    )
except ImportError as e:
    print(f"Dependency Error: {e}\n"
          f"Run: pip install flask opencv-python python-pptx yt-dlp pillow")
    sys.exit(1)

logging.basicConfig(level=logging.INFO,
                     format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

# Shared, on-disk preview folder next to this script. Wiped at the start
# of every new extraction job — it is a "live view" folder only.
BASE_DIR = Path(__file__).resolve().parent
PREVIEW_DIR = BASE_DIR / "preview_frames"


# --------------------------------------------------------------------------- #
#  Small helpers
# --------------------------------------------------------------------------- #

def clear_preview_dir() -> None:
    """Wipe and recreate the shared preview folder. Called once at the
    start of every new extraction job so old frames never linger."""
    shutil.rmtree(PREVIEW_DIR, ignore_errors=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)


def quality_label(width, height) -> str:
    if not width or not height:
        return "Unknown resolution"
    longest = max(width, height)
    if longest >= 3840:
        tag = "4K/8K Ultra HD"
    elif longest >= 2560:
        tag = "2K QHD"
    elif longest >= 1920:
        tag = "Full HD 1080p"
    elif longest >= 1280:
        tag = "HD 720p"
    else:
        tag = "SD"
    return f"{width}x{height} — {tag}"


def format_duration(seconds) -> str:
    if not seconds or seconds <= 0:
        return "unknown"
    seconds = int(seconds)
    h, rem = divmod(seconds, 3600)
    m, s = divmod(rem, 60)
    return f"{h:02d}:{m:02d}:{s:02d}" if h else f"{m:02d}:{s:02d}"


def add_picture_fit(slide, img_path, slide_w, slide_h):
    """Insert an image centered on the slide, scaled to fit WITHOUT
    stretching / distorting the original aspect ratio."""
    with PILImage.open(img_path) as im:
        iw, ih = im.size

    img_ratio = iw / ih
    slide_ratio = slide_w / slide_h

    if img_ratio > slide_ratio:
        w = slide_w
        h = int(slide_w / img_ratio)
    else:
        h = slide_h
        w = int(slide_h * img_ratio)

    left = int((slide_w - w) / 2)
    top = int((slide_h - h) / 2)
    slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)


# --------------------------------------------------------------------------- #
#  In-memory job registry (shared preview folder -> one active job at a time)
# --------------------------------------------------------------------------- #

jobs = {}            # job_id -> job state dict
jobs_lock = threading.Lock()
current_job_id = None  # the one job allowed to write into PREVIEW_DIR


class JobReporter:
    """Small helper that lets the background extraction thread push
    live log lines / progress / quality info into the shared job dict
    without every method needing to know about locking."""

    def __init__(self, job_id):
        self.job_id = job_id

    def log(self, msg: str) -> None:
        ts = datetime.now().strftime('%H:%M:%S')
        line = f"[{ts}] {msg}"
        with jobs_lock:
            jobs[self.job_id]["log"].append(line)
        logger.info(line)

    def progress(self, pct) -> None:
        with jobs_lock:
            jobs[self.job_id]["progress"] = max(0, min(100, int(pct)))

    def quality(self, label, duration) -> None:
        with jobs_lock:
            jobs[self.job_id]["quality_label"] = label
            jobs[self.job_id]["duration"] = duration
            jobs[self.job_id]["duration_label"] = format_duration(duration)


# --------------------------------------------------------------------------- #
#  Core conversion logic
# --------------------------------------------------------------------------- #

class VideoExtractor:
    """Handles frame extraction, live preview and PPTX/ZIP export for a
    single conversion job."""

    def __init__(self, source: str, interval: int = 1):
        self.source = source.strip()
        self.interval = max(1, int(interval))
        self.is_url = self.source.lower().startswith(('http://', 'https://', 'www.'))

        self.saved_frames = []   # list[Path] — inside PREVIEW_DIR
        self.base_name = "Online_Video"

        if not self.is_url:
            local_path = Path(self.source)
            if not local_path.exists():
                raise FileNotFoundError(f"Local file not found: {self.source}")
            self.base_name = local_path.stem

    def _resolve_stream(self, reporter: JobReporter):
        """Returns (capture_target, duration_seconds). Never downloads a
        remote source to disk — yt-dlp only resolves a direct stream URL,
        which OpenCV then reads from progressively."""
        if not self.is_url:
            reporter.log(f"Opening local file: {self.source}")
            return str(self.source), None

        reporter.log(f"Resolving best-quality stream for: {self.source}")
        ydl_opts = {
            # No 'ext=mp4' restriction — that filter silently drops
            # 1440p/4K/8K streams, which are almost always webm.
            'format': 'bestvideo/best',
            'format_sort': ['res', 'fps', 'vcodec:vp9.2', 'vcodec:vp9', 'br'],
            'quiet': True,
            'no_warnings': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(self.source, download=False)

        if info.get('title'):
            self.base_name = "".join(
                c for c in info['title'] if c.isalnum() or c in (' ', '_', '-')
            ).rstrip().replace(' ', '_') or self.base_name

        width, height = info.get('width'), info.get('height')
        duration = info.get('duration')
        label = quality_label(width, height)
        reporter.quality(label, duration)
        reporter.log(f"Detected quality: {label}")
        reporter.log(f"Video length: {format_duration(duration)}"
                      + (" (live/unknown length)" if not duration else ""))
        reporter.log("Opening the remote stream directly — the file is "
                      "NOT being downloaded to disk.")
        return info['url'], duration

    def extract_frames(self, reporter: JobReporter) -> None:
        target, duration = self._resolve_stream(reporter)

        cap = cv2.VideoCapture(target)
        if not cap.isOpened():
            raise ValueError("Could not open the video stream/source.")

        fps = cap.get(cv2.CAP_PROP_FPS)
        if fps <= 0:
            fps = 25  # fallback for streams that hide metadata

        if not self.is_url:
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            frame_total = cap.get(cv2.CAP_PROP_FRAME_COUNT)
            duration = (frame_total / fps) if frame_total and fps else None
            label = quality_label(width, height)
            reporter.quality(label, duration)
            reporter.log(f"Detected quality: {label}")
            reporter.log(f"Video length: {format_duration(duration)}")

        reporter.log(f"Extracting one frame every {self.interval}s ...")

        frame_gap = max(1, int(fps * self.interval))
        frame_idx = 0
        saved_idx = 0

        while True:
            success, frame = cap.read()
            if not success:
                break

            if frame_idx % frame_gap == 0:
                target_img_path = PREVIEW_DIR / f"slide_{saved_idx:04d}.jpg"
                # Max JPEG quality (OpenCV default is ~95, we force 100)
                cv2.imwrite(str(target_img_path), frame, [cv2.IMWRITE_JPEG_QUALITY, 100])
                self.saved_frames.append(target_img_path)

                t = frame_idx / fps
                reporter.log(f"Saved frame #{saved_idx + 1} (t={t:0.1f}s)")
                if duration:
                    reporter.progress((t / duration) * 100)
                saved_idx += 1

            frame_idx += 1

        cap.release()
        reporter.progress(100)
        reporter.log(f"Extraction complete — {len(self.saved_frames)} frames saved "
                      f"to {PREVIEW_DIR.name}/")

    def get_selected_frames(self, selected_indices):
        return [self.saved_frames[i] for i in selected_indices
                if 0 <= i < len(self.saved_frames)]

    def save_to_pptx(self, selected_indices) -> Path:
        frames = self.get_selected_frames(selected_indices)
        if not frames:
            raise ValueError("No frames selected for export.")

        prs = Presentation()
        prs.slide_width, prs.slide_height = SLIDE_WIDTH, SLIDE_HEIGHT

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Video Content Export"
        source_label = self.source if not self.is_url else "Remote Video Source"
        title_slide.placeholders[1].text = (
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"Source: {source_label}"
        )

        blank_layout = prs.slide_layouts[6]
        for img_path in frames:
            slide = prs.slides.add_slide(blank_layout)
            add_picture_fit(slide, img_path, prs.slide_width, prs.slide_height)

        export_dir = Path(tempfile.mkdtemp(prefix="export_"))
        output_path = export_dir / f"{self.base_name}_output.pptx"
        prs.save(str(output_path))
        logger.info(f"PPTX saved: {output_path}")
        return output_path

    def export_images(self, selected_indices) -> Path:
        """Zips the selected frames (copies of the preview files) and
        returns the zip path."""
        frames = self.get_selected_frames(selected_indices)
        if not frames:
            raise ValueError("No frames selected for export.")

        export_dir = Path(tempfile.mkdtemp(prefix="export_"))
        images_dir = export_dir / f"{self.base_name}_HD_Frames"
        images_dir.mkdir(exist_ok=True)
        for img_path in frames:
            shutil.copy(img_path, images_dir / img_path.name)

        zip_base = export_dir / f"{self.base_name}_HD_Frames"
        zip_path = shutil.make_archive(str(zip_base), 'zip', root_dir=str(images_dir))
        return Path(zip_path)


# --------------------------------------------------------------------------- #
#  Flask app
# --------------------------------------------------------------------------- #

app = Flask(__name__)


def run_extraction_job(job_id: str, extractor: VideoExtractor):
    reporter = JobReporter(job_id)
    try:
        extractor.extract_frames(reporter)
        with jobs_lock:
            jobs[job_id]["status"] = "done"
    except Exception as e:
        logger.error(f"Job {job_id} failed: {e}")
        reporter.log(f"ERROR: {e}")
        with jobs_lock:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = str(e)


@app.route('/')
def index():
    return render_template_string(INDEX_HTML)


@app.route('/extract', methods=['POST'])
def extract():
    global current_job_id

    data = request.get_json(force=True)
    source = (data.get('source') or '').strip()
    interval = data.get('interval', 1)

    if not source:
        return jsonify({"error": "Please provide a video URL or local file path."}), 400

    with jobs_lock:
        if current_job_id and jobs.get(current_job_id, {}).get("status") == "processing":
            return jsonify({"error": "An extraction is already running. "
                                      "Please wait for it to finish."}), 409

    try:
        extractor = VideoExtractor(source, interval=interval)
    except FileNotFoundError as e:
        return jsonify({"error": str(e)}), 400

    # Wipe the shared preview folder before this job starts writing to it.
    clear_preview_dir()

    job_id = str(uuid.uuid4())
    with jobs_lock:
        jobs[job_id] = {
            "status": "processing",
            "extractor": extractor,
            "error": None,
            "log": [],
            "progress": 0,
            "quality_label": None,
            "duration": None,
            "duration_label": None,
        }
        current_job_id = job_id

    threading.Thread(target=run_extraction_job, args=(job_id, extractor), daemon=True).start()
    return jsonify({"job_id": job_id, "preview_folder": str(PREVIEW_DIR)})


@app.route('/status/<job_id>')
def status(job_id):
    since = request.args.get('since', 0, type=int)
    with jobs_lock:
        job = jobs.get(job_id)
        if not job:
            return jsonify({"status": "not_found"}), 404
        logs = job["log"][since:]
        return jsonify({
            "status": job["status"],
            "error": job["error"],
            "frame_count": len(job["extractor"].saved_frames),
            "progress": job["progress"],
            "quality_label": job["quality_label"],
            "duration_label": job["duration_label"],
            "logs": logs,
            "log_cursor": len(job["log"]),
        })


@app.route('/frame/<job_id>/<int:idx>')
def frame(job_id, idx):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        abort(404)
    frames = job["extractor"].saved_frames
    if idx < 0 or idx >= len(frames):
        abort(404)
    return send_file(str(frames[idx]), mimetype='image/jpeg')


@app.route('/export/pptx/<job_id>', methods=['POST'])
def export_pptx(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        abort(404)
    data = request.get_json(force=True)
    selected = data.get('selected', [])
    try:
        output_path = job["extractor"].save_to_pptx(selected)
    except ValueError as e:
        return jsonify({"error": str(e)}), 400

    schedule_temp_cleanup(output_path.parent)
    return send_file(str(output_path), as_attachment=True,
                      download_name=output_path.name)


@app.route('/export/images/<job_id>', methods=['POST'])
def export_images(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        abort(404)
    data = request.get_json(force=True)
    selected = data.get('selected', [])
    try:
        zip_path = job["extractor"].export_images(selected)
    except ValueError as e:
        return jsonify({"error": str(e)}), 400

    schedule_temp_cleanup(zip_path.parent)
    return send_file(str(zip_path), as_attachment=True, download_name=zip_path.name)


def schedule_temp_cleanup(dir_to_remove: Path):
    """Deletes a one-off export temp dir (pptx/zip) once the response
    for this request has been sent. Independent of the shared preview
    folder, which is only cleared at the start of the next extraction."""
    @app.after_this_request
    def _cleanup(response):
        shutil.rmtree(dir_to_remove, ignore_errors=True)
        return response


@app.route('/cleanup/<job_id>', methods=['POST'])
def cleanup(job_id):
    global current_job_id
    with jobs_lock:
        jobs.pop(job_id, None)
        if current_job_id == job_id:
            current_job_id = None
    clear_preview_dir()
    return jsonify({"ok": True})


# --------------------------------------------------------------------------- #
#  Inline HTML/JS UI
# --------------------------------------------------------------------------- #

INDEX_HTML = """
<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Video to Presentation Converter</title>
<style>
  * { box-sizing: border-box; }
  body {
    font-family: 'Segoe UI', Arial, sans-serif;
    margin: 0; padding: 30px 15px 60px;
    min-height: 100vh;
    background: linear-gradient(135deg, #1f1147 0%, #3a1c71 35%, #6a3093 65%, #a44a9c 100%);
    color: #222;
  }
  .wrap { max-width: 1000px; margin: 0 auto; }
  h1 {
    color: #fff; text-align: center; font-size: 26px; margin-bottom: 4px;
    text-shadow: 0 2px 6px rgba(0,0,0,0.3);
  }
  .subtitle { text-align: center; color: rgba(255,255,255,0.75); margin-bottom: 24px; font-size: 13px; }

  .card {
    background: rgba(255,255,255,0.95);
    border-radius: 14px;
    padding: 20px 22px;
    margin-bottom: 18px;
    box-shadow: 0 8px 30px rgba(0,0,0,0.25);
    backdrop-filter: blur(6px);
  }

  .row { display: flex; gap: 10px; flex-wrap: wrap; }
  input[type=text] {
    flex: 1; min-width: 200px; padding: 11px 14px; border: 1px solid #ddd;
    border-radius: 8px; font-size: 14px; outline: none;
  }
  input[type=text]:focus { border-color: #8e44ad; }
  input[type=number] { width: 90px; padding: 11px; border: 1px solid #ddd; border-radius: 8px; font-size: 14px; }
  button {
    padding: 11px 20px; cursor: pointer; border: none; border-radius: 8px;
    font-size: 14px; font-weight: 600; color: #fff;
    background: linear-gradient(135deg, #8e44ad, #3a1c71);
    transition: transform .12s ease, box-shadow .12s ease;
  }
  button:hover { transform: translateY(-1px); box-shadow: 0 4px 14px rgba(58,28,113,0.4); }
  button:disabled { opacity: .5; cursor: not-allowed; transform: none; box-shadow: none; }
  button.secondary { background: linear-gradient(135deg, #666, #333); }

  .meta-line { font-size: 13px; color: #555; margin-top: 10px; }
  .badge {
    display: inline-block; padding: 3px 10px; border-radius: 20px; font-size: 12px;
    font-weight: 700; color: #fff; background: linear-gradient(135deg, #16a085, #2ecc71);
    margin-left: 6px;
  }

  .spinner {
    display: inline-block; width: 14px; height: 14px; border-radius: 50%;
    border: 2px solid rgba(142,68,173,0.25); border-top-color: #8e44ad;
    animation: spin .8s linear infinite; margin-right: 8px; vertical-align: -2px;
  }
  @keyframes spin { to { transform: rotate(360deg); } }

  .progress-outer {
    background: #eee; border-radius: 20px; height: 16px; overflow: hidden; margin-top: 12px;
  }
  .progress-inner {
    height: 100%; width: 0%; border-radius: 20px;
    background: linear-gradient(90deg, #8e44ad, #e056fd, #8e44ad);
    background-size: 200% 100%;
    animation: shimmer 1.4s linear infinite;
    transition: width .35s ease;
  }
  @keyframes shimmer { from { background-position: 0% 0; } to { background-position: 200% 0; } }
  .progress-text { font-size: 12px; color: #555; margin-top: 4px; text-align: right; }

  #console {
    background: #12121c; color: #b9f6ca; font-family: 'Consolas', monospace;
    font-size: 12.5px; border-radius: 8px; padding: 12px 14px;
    height: 160px; overflow-y: auto; margin-top: 14px; line-height: 1.5;
  }
  #console div.errline { color: #ff7675; }

  #grid {
    display: grid; grid-template-columns: repeat(auto-fill, minmax(140px, 1fr));
    gap: 12px; margin-top: 6px;
  }
  .cell {
    position: relative; border-radius: 10px; overflow: hidden; cursor: pointer;
    box-shadow: 0 3px 10px rgba(0,0,0,0.15);
    animation: fadeInUp .35s ease both;
  }
  @keyframes fadeInUp { from { opacity: 0; transform: translateY(8px); } to { opacity: 1; transform: translateY(0); } }
  .cell img { width: 100%; display: block; aspect-ratio: 16/9; object-fit: cover; }
  .cell .idx {
    position: absolute; bottom: 4px; right: 6px; color: #fff; font-size: 11px;
    background: rgba(0,0,0,0.55); padding: 1px 6px; border-radius: 10px;
  }
  .cell .chk {
    position: absolute; top: 6px; left: 6px; width: 18px; height: 18px; z-index: 2;
  }

  .actions { display: flex; gap: 10px; margin-top: 16px; flex-wrap: wrap; }
  .actions button { flex: 1; min-width: 160px; }

  #lightbox {
    position: fixed; inset: 0; background: rgba(0,0,0,0.9); display: none;
    align-items: center; justify-content: center; z-index: 999; flex-direction: column;
  }
  #lightbox img { max-width: 90vw; max-height: 78vh; border-radius: 8px; box-shadow: 0 10px 40px rgba(0,0,0,0.6); }
  #lightbox .lb-counter { color: #fff; margin-top: 10px; font-size: 13px; opacity: .85; }
  .lb-btn {
    position: absolute; top: 50%; transform: translateY(-50%);
    background: rgba(255,255,255,0.15); color: #fff; border: none; font-size: 26px;
    width: 46px; height: 46px; border-radius: 50%; cursor: pointer;
  }
  .lb-btn:hover { background: rgba(255,255,255,0.3); }
  #lb-prev { left: 18px; } #lb-next { right: 18px; }
  #lb-close {
    position: absolute; top: 18px; right: 22px; background: transparent; color: #fff;
    font-size: 26px; border: none; cursor: pointer;
  }
</style>
</head>
<body>
<div class="wrap">
  <h1>🎬 Video → Presentation Converter</h1>
  <div class="subtitle">Extract frames, preview them live, and export to PPTX or images</div>

  <div class="card">
    <div class="row">
      <input type="text" id="source" placeholder="Video URL (YouTube etc.) or local file path">
      <input type="number" id="interval" value="1" min="1" title="Seconds between frames">
      <button id="extractBtn" onclick="startExtraction()">Extract Frames</button>
    </div>
    <div class="meta-line" id="qualityLine"></div>
    <div class="meta-line" id="folderLine" style="color:#8e44ad;"></div>
  </div>

  <div class="card" id="progressCard" style="display:none;">
    <div><span class="spinner" id="spinnerIcon"></span><span id="statusText">Starting…</span></div>
    <div class="progress-outer"><div class="progress-inner" id="progressBar"></div></div>
    <div class="progress-text" id="progressPct">0%</div>
    <div id="console"></div>
  </div>

  <div class="card" id="gridCard" style="display:none;">
    <div class="row" style="margin-bottom:12px;">
      <button class="secondary" onclick="selectAll(true)">Select All</button>
      <button class="secondary" onclick="selectAll(false)">Deselect All</button>
    </div>
    <div id="grid"></div>
    <div class="actions">
      <button onclick="exportPptx()">📊 Save Selected as PPTX</button>
      <button onclick="exportImages()">🖼️ Save Selected Images (.zip)</button>
    </div>
  </div>
</div>

<div id="lightbox">
  <button id="lb-close" onclick="closeLightbox()">✕</button>
  <button id="lb-prev" class="lb-btn" onclick="lbMove(-1)">‹</button>
  <img id="lb-img" src="">
  <button id="lb-next" class="lb-btn" onclick="lbMove(1)">›</button>
  <div class="lb-counter" id="lb-counter"></div>
</div>

<script>
let jobId = null;
let frameShown = 0;
let frameCount = 0;
let logCursor = 0;
let pollTimer = null;
let lbIndex = 0;

function el(id) { return document.getElementById(id); }

function resetUI() {
  el('grid').innerHTML = '';
  el('console').innerHTML = '';
  el('gridCard').style.display = 'none';
  el('progressCard').style.display = 'block';
  el('qualityLine').innerHTML = '';
  el('progressBar').style.width = '0%';
  el('progressPct').innerText = '0%';
  el('spinnerIcon').style.display = 'inline-block';
  frameShown = 0; frameCount = 0; logCursor = 0;
}

async function startExtraction() {
  const source = el('source').value.trim();
  const interval = parseInt(el('interval').value || '1', 10);
  if (!source) { alert('Please provide a video URL or local file path.'); return; }

  resetUI();
  el('statusText').innerText = 'Starting extraction…';
  el('extractBtn').disabled = true;

  const res = await fetch('/extract', {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({source, interval})
  });
  const data = await res.json();
  if (data.error) {
    el('statusText').innerText = 'Error: ' + data.error;
    el('spinnerIcon').style.display = 'none';
    el('extractBtn').disabled = false;
    return;
  }

  jobId = data.job_id;
  el('folderLine').innerText = '📁 Live preview folder: ' + data.preview_folder;
  pollTimer = setInterval(poll, 700);
}

async function poll() {
  const res = await fetch(`/status/${jobId}?since=${logCursor}`);
  const data = await res.json();

  if (data.logs && data.logs.length) {
    const c = el('console');
    data.logs.forEach(line => {
      const d = document.createElement('div');
      if (/ERROR/i.test(line)) d.className = 'errline';
      d.innerText = line;
      c.appendChild(d);
    });
    c.scrollTop = c.scrollHeight;
    logCursor = data.log_cursor;
  }

  if (data.quality_label) {
    el('qualityLine').innerHTML = 'Quality detected: <span class="badge">' + data.quality_label + '</span>'
      + (data.duration_label ? '  ·  Length: ' + data.duration_label : '');
  }

  el('progressBar').style.width = (data.progress || 0) + '%';
  el('progressPct').innerText = (data.progress || 0) + '%';

  frameCount = data.frame_count;
  if (frameCount > frameShown) growGrid();

  if (data.status === 'processing') {
    el('statusText').innerText = `Extracting… ${frameCount} frame(s) so far`;
  } else if (data.status === 'done') {
    clearInterval(pollTimer);
    el('statusText').innerText = `Done — ${frameCount} frames extracted`;
    el('spinnerIcon').style.display = 'none';
    el('extractBtn').disabled = false;
    el('gridCard').style.display = 'block';
  } else if (data.status === 'error') {
    clearInterval(pollTimer);
    el('statusText').innerText = 'Error: ' + data.error;
    el('spinnerIcon').style.display = 'none';
    el('extractBtn').disabled = false;
  }
}

function growGrid() {
  const grid = el('grid');
  for (let i = frameShown; i < frameCount; i++) {
    const cell = document.createElement('div');
    cell.className = 'cell';
    cell.innerHTML = `
      <input type="checkbox" class="frame-check chk" data-idx="${i}" checked onclick="event.stopPropagation()">
      <img src="/frame/${jobId}/${i}" loading="lazy" onclick="openLightbox(${i})">
      <div class="idx">#${i + 1}</div>
    `;
    grid.appendChild(cell);
  }
  frameShown = frameCount;
}

function selectAll(state) {
  document.querySelectorAll('.frame-check').forEach(cb => cb.checked = state);
}

function getSelected() {
  return Array.from(document.querySelectorAll('.frame-check'))
    .filter(cb => cb.checked)
    .map(cb => parseInt(cb.dataset.idx, 10));
}

function openLightbox(idx) {
  lbIndex = idx;
  renderLightbox();
  el('lightbox').style.display = 'flex';
}
function closeLightbox() { el('lightbox').style.display = 'none'; }
function lbMove(delta) {
  lbIndex = (lbIndex + delta + frameCount) % frameCount;
  renderLightbox();
}
function renderLightbox() {
  el('lb-img').src = `/frame/${jobId}/${lbIndex}`;
  el('lb-counter').innerText = `${lbIndex + 1} / ${frameCount}`;
}
document.addEventListener('keydown', (e) => {
  if (el('lightbox').style.display !== 'flex') return;
  if (e.key === 'Escape') closeLightbox();
  if (e.key === 'ArrowLeft') lbMove(-1);
  if (e.key === 'ArrowRight') lbMove(1);
});
el('lightbox') && el('lightbox').addEventListener('click', (e) => {
  if (e.target.id === 'lightbox') closeLightbox();
});

async function exportPptx() {
  const selected = getSelected();
  if (!selected.length) { alert('Select at least one frame.'); return; }
  el('statusText').innerText = 'Building PPTX…';
  const res = await fetch('/export/pptx/' + jobId, {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({selected})
  });
  if (!res.ok) { const err = await res.json(); alert('Error: ' + err.error); return; }
  const blob = await res.blob();
  downloadBlob(blob, 'presentation.pptx');
}

async function exportImages() {
  const selected = getSelected();
  if (!selected.length) { alert('Select at least one frame.'); return; }
  el('statusText').innerText = 'Zipping images…';
  const res = await fetch('/export/images/' + jobId, {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({selected})
  });
  if (!res.ok) { const err = await res.json(); alert('Error: ' + err.error); return; }
  const blob = await res.blob();
  downloadBlob(blob, 'frames.zip');
}

function downloadBlob(blob, filename) {
  const url = URL.createObjectURL(blob);
  const a = document.createElement('a');
  a.href = url; a.download = filename;
  document.body.appendChild(a); a.click(); a.remove();
  URL.revokeObjectURL(url);
}
</script>
</body>
</html>
"""


if __name__ == '__main__':
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    app.run(debug=True, host='127.0.0.1', port=8080)